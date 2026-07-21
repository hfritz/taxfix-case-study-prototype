"""
Builds process-case-study-deck.pptx — a real PowerPoint/Keynote/Google
Slides file, using the same content and brand system (colors from
brand-brief.md, Archivo/Archivo Black type) as the HTML version.

Font note: Archivo / Archivo Black are specified by name, same fallback
approach the brand kit's own template.html uses — if the viewing machine
doesn't have them installed, PowerPoint/Keynote substitute a system
sans-serif automatically. python-pptx cannot reliably embed fonts across
platforms, so this matches the brand kit's own documented approach
rather than being a compromise unique to this file.

Re-run this file (`python3 build_pptx.py`) after editing to regenerate
process-case-study-deck.pptx — easier than hand-editing the binary file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------------------------------------------------------------
# BRAND TOKENS — copied verbatim from brand-brief.md
# ---------------------------------------------------------------
EG_LIGHT = RGBColor.from_string("ECFFC7")
EG_CALM = RGBColor.from_string("CEF5A4")
EG_VIVID = RGBColor.from_string("ADEE68")
EG_DARK = RGBColor.from_string("36893B")
EG_VERYDARK = RGBColor.from_string("154618")

OW_LIGHT = RGBColor.from_string("FDF8F2")
OW_CALM = RGBColor.from_string("EAE0D7")
OW_VIVID = RGBColor.from_string("9A9288")

GOLD_VIVID = RGBColor.from_string("F8A21A")
GOLD_DARK = RGBColor.from_string("2B1A00")
BLUE_VIVID = RGBColor.from_string("668CFF")

WHITE = RGBColor.from_string("FFFFFF")
BLACK = RGBColor.from_string("000000")

FONT_HEAD = "Archivo Black"
FONT_BODY = "Archivo"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------
# HELPERS
# ---------------------------------------------------------------
def add_slide(bg_color):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    return slide


def add_text(slide, left, top, width, height, text, *, font=FONT_BODY, size=14,
             color=BLACK, bold=False, italic=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_runs(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.3, wrap=True):
    """runs: list of dicts {text, font, size, color, bold, italic}"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in runs:
        run = p.add_run()
        run.text = r["text"]
        run.font.name = r.get("font", FONT_BODY)
        run.font.size = Pt(r.get("size", 14))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", BLACK)
    return box


def add_card(slide, left, top, width, height, fill, *, line_color=None,
             corner=0.06):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = corner
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_dashed_card(slide, left, top, width, height, line_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.08
    shp.fill.background()
    shp.line.color.rgb = line_color
    shp.line.width = Pt(1)
    shp.line.dash_style = 4  # MSO_LINE_DASH_STYLE.DASH
    shp.shadow.inherit = False
    return shp


def add_pill(slide, left, top, width, height, text, fill, text_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT_HEAD
    run.font.size = Pt(9)
    run.font.color.rgb = text_color
    return shp


def add_dot(slide, cx, cy, r, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_motif(slide, left, top, size, color):
    """The deconstructed percent sign: dot, diagonal slash, dot."""
    r = size * 0.09
    add_dot(slide, left + size * 0.18, top + size * 0.18, r, color)
    line = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Emu(int(left + size * 0.30)), Emu(int(top + size * 0.30)),
        Emu(int(size * 0.5)), Emu(int(size * 0.12)),
    )
    line.rotation = -40
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    line.adjustments[0] = 0.5
    add_dot(slide, left + size * 0.82, top + size * 0.82, r, color)


def footer(slide, index, total, color):
    add_text(
        slide, MARGIN, SLIDE_H - Inches(0.55), Inches(2), Inches(0.35),
        f"{index:02d} / {total:02d}", font=FONT_HEAD, size=11, color=color,
    )


def eyebrow(slide, text, color, *, left=MARGIN, top=Inches(0.7), width=Inches(8),
            align=PP_ALIGN.LEFT):
    add_text(slide, left, top, width, Inches(0.35), text, font=FONT_HEAD,
              size=12, color=color, align=align)


TOTAL = 10

# =================================================================
# SLIDE 1 — COVER
# =================================================================
s = add_slide(EG_VERYDARK)
add_motif(s, Inches(8.6), Inches(2.6), Inches(3.6), EG_DARK)
eyebrow(s, "AI PM BUILDER (ASSISTED) — CASE STUDY", EG_CALM, top=Inches(1.0))
add_text(
    s, MARGIN, Inches(1.5), Inches(8.2), Inches(2.4),
    "A premium Expert Service for the freelancers Taxfix says no to.",
    font=FONT_HEAD, size=40, color=WHITE, line_spacing=1.05,
)
add_text(
    s, MARGIN, Inches(4.0), Inches(7.2), Inches(0.9),
    "How I approached the segment, the price, the build, and the AI "
    "workflow behind all three.",
    font=FONT_BODY, size=16, color=OW_LIGHT,
)
tile = add_card(s, MARGIN, Inches(5.15), Inches(0.42), Inches(0.42), EG_VIVID, corner=0.3)
tile.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tile.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "·/·"
r.font.name = FONT_HEAD
r.font.size = Pt(13)
r.font.color.rgb = EG_VERYDARK
add_text(
    s, MARGIN + Inches(0.58), Inches(5.22), Inches(7), Inches(0.4),
    "Helmut Fritz · 2026 · Unofficial prototype, not affiliated "
    "with or endorsed by Taxfix",
    font=FONT_BODY, size=11, color=RGBColor.from_string("D8D3C9"),
)
footer(s, 1, TOTAL, OW_LIGHT)

# =================================================================
# SLIDE 2 — THE SEGMENT
# =================================================================
s = add_slide(OW_LIGHT)
eyebrow(s, "01 · THE SEGMENT", EG_DARK)
add_text(
    s, MARGIN, Inches(1.15), Inches(8.5), Inches(1.6),
    "Not “self-employed.” This exact intersection.",
    font=FONT_HEAD, size=30, color=EG_VERYDARK, line_spacing=1.05,
)
add_text(
    s, MARGIN, Inches(2.75), Inches(8.2), Inches(1.3),
    "Lena: a regelbesteuerte freelancer past the Kleinunternehmer line, "
    "with a client or income stream outside Germany. Either trait alone "
    "is already served today. Both at once is where Taxfix’s product "
    "stops and a slow, expensive Steuerberater takes over.",
    font=FONT_BODY, size=13.5, color=RGBColor.from_string("2E4A30"),
)
card_w = Inches(3.6)
c1 = add_card(s, MARGIN, Inches(4.4), card_w, Inches(1.3), WHITE, line_color=OW_CALM)
c1.text_frame.margin_left = Inches(0.22)
c1.text_frame.margin_top = Inches(0.18)
c1.text_frame.margin_right = Inches(0.2)
p = c1.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Self-employed"; r.font.name = FONT_HEAD; r.font.size = Pt(14); r.font.color.rgb = EG_VERYDARK
p2 = c1.text_frame.add_paragraph()
p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "Past the Kleinunternehmer line — regelbesteuert, not the small-business exemption."
r2.font.name = FONT_BODY; r2.font.size = Pt(11); r2.font.color.rgb = RGBColor.from_string("2E4A30")
c1.text_frame.word_wrap = True

add_text(s, MARGIN + card_w + Inches(0.15), Inches(4.85), Inches(0.4), Inches(0.5), "+",
          font=FONT_HEAD, size=22, color=EG_DARK, align=PP_ALIGN.CENTER)

c2 = add_card(s, MARGIN + card_w + Inches(0.55), Inches(4.4), card_w, Inches(1.3), WHITE, line_color=OW_CALM)
c2.text_frame.margin_left = Inches(0.22)
c2.text_frame.margin_top = Inches(0.18)
c2.text_frame.margin_right = Inches(0.2)
p = c2.text_frame.paragraphs[0]
r = p.add_run(); r.text = "Cross-border"; r.font.name = FONT_HEAD; r.font.size = Pt(14); r.font.color.rgb = EG_VERYDARK
p2 = c2.text_frame.add_paragraph()
p2.space_before = Pt(6)
r2 = p2.add_run(); r2.text = "A foreign client, foreign rental income, or a double-taxation-treaty case."
r2.font.name = FONT_BODY; r2.font.size = Pt(11); r2.font.color.rgb = RGBColor.from_string("2E4A30")
c2.text_frame.word_wrap = True
footer(s, 2, TOTAL, EG_VERYDARK)

# =================================================================
# SLIDE 3 — THE GAP
# =================================================================
s = add_slide(OW_CALM)
eyebrow(s, "02 · THE GAP", EG_DARK)
add_text(
    s, MARGIN, Inches(1.15), Inches(9.5), Inches(1.3),
    "Three weaknesses. Two are direct quotes, one is a judgment call.",
    font=FONT_HEAD, size=27, color=EG_VERYDARK, line_spacing=1.05,
)

gap_rows = [
    ("QUOTE", EG_VIVID, EG_VERYDARK,
     "“Gewerbetreibende mit Regelversteuerung” — named on Taxfix’s "
     "own Experten-Service page as not supported."),
    ("QUOTE", EG_VIVID, EG_VERYDARK,
     "“Einige internationale Steuerfälle” and Auslandsvermietung — "
     "also named exclusions, not an inference."),
    ("INFERENCE", GOLD_VIVID, GOLD_DARK,
     "One annual filing only, not her real quarterly cadence. My own "
     "scoping read, not a documented gap — flagged, and kept out of "
     "v1 on purpose."),
]
y = Inches(2.75)
row_h = Inches(1.15)
for tag, tag_fill, tag_color, body in gap_rows:
    add_card(s, MARGIN, y, Inches(10.5), row_h, WHITE, line_color=OW_VIVID)
    add_pill(s, MARGIN + Inches(0.25), y + Inches(0.22), Inches(1.15), Inches(0.32), tag, tag_fill, tag_color)
    add_text(s, MARGIN + Inches(1.6), y + Inches(0.18), Inches(8.5), Inches(0.85), body,
              font=FONT_BODY, size=12.5, color=EG_VERYDARK, anchor=MSO_ANCHOR.MIDDLE)
    y = y + row_h + Inches(0.2)
footer(s, 3, TOTAL, EG_VERYDARK)

# =================================================================
# SLIDE 4 — THE PRICE
# =================================================================
s = add_slide(EG_VIVID)
eyebrow(s, "03 · THE PRICE", RGBColor.from_string("3F6B26"), left=Inches(0), width=SLIDE_W, align=PP_ALIGN.CENTER)
add_text(s, Inches(0), Inches(1.15), SLIDE_W, Inches(2.0), "€449",
          font=FONT_HEAD, size=110, color=EG_VERYDARK, align=PP_ALIGN.CENTER)
add_text(
    s, Inches(1.8), Inches(3.35), Inches(9.7), Inches(0.7),
    "Flat. 5 advisor hours included. A disclosed surcharge beyond that "
    "— never a surprise bill.",
    font=FONT_BODY, size=15, color=EG_VERYDARK, align=PP_ALIGN.CENTER,
)

cards = [
    ("Doing it yourself", "6–10 hours of your own time — worth €480–800 "
     "at €80/hr, plus the risk of a mistake.", RGBColor.from_string("C3E796"), EG_VERYDARK),
    ("A traditional Kanzlei", "€520–1,500+ for a complex case, billed hourly. "
     "You don’t know the total until it arrives.", RGBColor.from_string("C3E796"), EG_VERYDARK),
    ("This tier", "€449, decided up front. Known before you start, not after.",
     EG_VERYDARK, EG_LIGHT),
]
cw = Inches(3.35)
gap = Inches(0.2)
start_x = (SLIDE_W - (cw * 3 + gap * 2)) / 2
x = start_x
for title, body, fill, color in cards:
    add_card(s, x, Inches(4.35), cw, Inches(1.55), fill)
    add_text(s, x + Inches(0.22), Inches(4.52), cw - Inches(0.44), Inches(0.35), title,
              font=FONT_HEAD, size=13, color=color)
    add_text(s, x + Inches(0.22), Inches(4.9), cw - Inches(0.44), Inches(0.9), body,
              font=FONT_BODY, size=10.5, color=color)
    x = x + cw + gap
footer(s, 4, TOTAL, RGBColor.from_string("3F6B26"))

# =================================================================
# SLIDE 5 — WHAT'S IN / OUT
# =================================================================
s = add_slide(OW_LIGHT)
eyebrow(s, "04 · WHAT'S IN, WHAT'S OUT", EG_DARK)
add_text(
    s, MARGIN, Inches(1.15), Inches(9.5), Inches(1.3),
    "A complexity-based tier, not a “more of everything” one.",
    font=FONT_HEAD, size=27, color=EG_VERYDARK, line_spacing=1.05,
)

included = [
    "A licensed Steuerberater with real self-employed + cross-border "
    "experience — not the generalist pool.",
    "A structured intake that sizes complexity before any work starts.",
    "One round of advisor follow-up and review before submission.",
]
excluded = [
    "Live or phone consultations — stays async to protect margin.",
    "Quarterly VAT and bookkeeping — a different, bigger product. A "
    "roadmap note, not v1.",
    "Multi-year retroactive filings.",
]

def bullet_list(slide, left, top, width, label, label_color, items, dot_color, text_color):
    add_text(slide, left, top, width, Inches(0.3), label, font=FONT_HEAD, size=11.5, color=label_color)
    y = top + Inches(0.45)
    for item in items:
        add_dot(slide, left + Inches(0.06), y + Inches(0.1), Inches(0.045), dot_color)
        add_text(slide, left + Inches(0.22), y, width - Inches(0.22), Inches(0.55), item,
                  font=FONT_BODY, size=12, color=text_color)
        y = y + Inches(0.62)

bullet_list(s, MARGIN, Inches(2.8), Inches(5.0), "INCLUDED", EG_DARK, included, EG_VIVID, EG_VERYDARK)
bullet_list(s, MARGIN + Inches(5.4), Inches(2.8), Inches(5.0), "DELIBERATELY EXCLUDED",
             RGBColor.from_string("7A7268"), excluded, OW_VIVID, RGBColor.from_string("6B6459"))

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(5.75), Inches(0.04), Inches(0.5))
bar.fill.solid(); bar.fill.fore_color.rgb = EG_VIVID; bar.line.fill.background(); bar.shadow.inherit = False
add_text(s, MARGIN + Inches(0.2), Inches(5.75), Inches(7), Inches(0.5),
          "The lever is complexity, priced honestly — not a bigger flat tier.",
          font=FONT_BODY, size=12, italic=True, color=RGBColor.from_string("6B6459"))
footer(s, 5, TOTAL, EG_VERYDARK)

# =================================================================
# SLIDE 6 — THE UNIT ECONOMICS
# =================================================================
s = add_slide(EG_VERYDARK)
eyebrow(s, "05 · THE UNIT ECONOMICS", EG_CALM)
add_text(s, MARGIN, Inches(1.15), Inches(9), Inches(1.0), "What €449 has to cover.",
          font=FONT_HEAD, size=27, color=WHITE)

add_runs(
    s, MARGIN, Inches(2.35), Inches(10.5), Inches(0.6),
    [
        {"text": "€449  ", "font": FONT_HEAD, "size": 22, "color": EG_LIGHT},
        {"text": "−  ", "font": FONT_HEAD, "size": 18, "color": RGBColor.from_string("6E7E70")},
        {"text": "€60/hr × 5hrs  ", "font": FONT_HEAD, "size": 22, "color": EG_LIGHT},
        {"text": "−  ", "font": FONT_HEAD, "size": 18, "color": RGBColor.from_string("6E7E70")},
        {"text": "€20", "font": FONT_HEAD, "size": 22, "color": EG_LIGHT},
    ],
)
add_runs(
    s, MARGIN, Inches(2.85), Inches(10.5), Inches(0.3),
    [
        {"text": "price", "font": FONT_BODY, "size": 10, "color": RGBColor.from_string("A8AFA0")},
        {"text": "                                    advisor cost", "font": FONT_BODY, "size": 10, "color": RGBColor.from_string("A8AFA0")},
        {"text": "                    platform opex", "font": FONT_BODY, "size": 10, "color": RGBColor.from_string("A8AFA0")},
    ],
)

add_text(s, MARGIN, Inches(3.5), Inches(3.2), Inches(1.3), "28.7%",
          font=FONT_HEAD, size=64, color=GOLD_VIVID)
add_text(
    s, MARGIN + Inches(3.3), Inches(3.75), Inches(6.5), Inches(1.0),
    "gross margin at the base case — honest, and below the 40%+ bar "
    "this needs to reach once real hours-per-case data replaces the "
    "placeholder.",
    font=FONT_BODY, size=12, color=RGBColor.from_string("D8DED2"),
)

line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(5.15), Inches(11.0), Pt(1))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor.from_string("2E4A30"); line.line.fill.background(); line.shadow.inherit = False

add_runs(
    s, MARGIN, Inches(5.35), Inches(11.0), Inches(1.6),
    [
        {"text": "€349 only cleared 8.3% margin — too thin, so price moved to €449.  ",
         "font": FONT_BODY, "size": 13, "bold": True, "color": WHITE},
        {"text": "The real risk isn't the price: it's hours-per-case, still a "
                  "placeholder. Sensitivity analysis shows it swings margin more "
                  "than advisor rate and opex combined (~€300, vs. ~€250 and ~€20).",
         "font": FONT_BODY, "size": 13, "color": RGBColor.from_string("D8DED2")},
    ],
    line_spacing=1.35,
)
footer(s, 6, TOTAL, RGBColor.from_string("6E7E70"))

# =================================================================
# SLIDE 7 — THE ARTIFACT
# =================================================================
s = add_slide(EG_VIVID)
eyebrow(s, "06 · THE ARTIFACT", RGBColor.from_string("3F6B26"))
add_text(s, MARGIN, Inches(1.15), Inches(9), Inches(1.3),
          "Three live pages, not a deck about pages.",
          font=FONT_HEAD, size=28, color=EG_VERYDARK, line_spacing=1.05)

built = [
    ("01", "Homepage clone", "Real taxfix.de HTML, with a genuine third pricing card added."),
    ("02", "Premium landing page", "Positioning, price/package, and a real comparison table against the standard tier."),
    ("03", "This process page", "Persona reasoning, a live pricing-logic slider, and the actual prompts used — not smoothed over."),
]
y = Inches(2.85)
row_h = Inches(0.85)
for n, title, body in built:
    add_card(s, MARGIN, y, Inches(9.8), row_h, RGBColor.from_string("C3E796"))
    add_text(s, MARGIN + Inches(0.25), y + Inches(0.12), Inches(0.6), Inches(0.6), n,
              font=FONT_HEAD, size=13, color=EG_VERYDARK)
    add_text(s, MARGIN + Inches(0.95), y + Inches(0.1), Inches(8.4), Inches(0.32), title,
              font=FONT_HEAD, size=13.5, color=EG_VERYDARK)
    add_text(s, MARGIN + Inches(0.95), y + Inches(0.42), Inches(8.4), Inches(0.38), body,
              font=FONT_BODY, size=11, color=EG_VERYDARK)
    y = y + row_h + Inches(0.18)

add_text(s, MARGIN, y + Inches(0.1), Inches(9), Inches(0.5),
          "Built solo, end to end: research, the pricing model, brand extraction, specs, and code.",
          font=FONT_BODY, size=11.5, color=RGBColor.from_string("3F6B26"))
footer(s, 7, TOTAL, RGBColor.from_string("3F6B26"))

# =================================================================
# SLIDE 8 — THE AI WORKFLOW
# =================================================================
s = add_slide(EG_VERYDARK)
add_motif(s, Inches(9.8), Inches(4.9), Inches(2.4), BLUE_VIVID)
eyebrow(s, "07 · THE AI WORKFLOW", BLUE_VIVID)
add_text(s, MARGIN, Inches(1.15), Inches(8), Inches(1.3),
          "Research, cross-checked. Not research, trusted.",
          font=FONT_HEAD, size=28, color=WHITE, line_spacing=1.05)
add_text(s, MARGIN, Inches(2.55), Inches(9.5), Inches(0.5),
          "Every model output in this deck was re-run somewhere else before being trusted.",
          font=FONT_BODY, size=13.5, color=RGBColor.from_string("D8DED2"))

ai_rows = [
    ("NotebookLM ", "ran a 6-phase interview-coach prompt against the real case brief before any solution was proposed."),
    ("The persona ", "was defined once, then re-run independently in a second Claude session as a sanity check."),
    ("The price ", "came from a Python sensitivity model, not a guess — and its first answer (€349) was corrected, not accepted."),
    ("Claude Code ", "built all three pages end to end from the resulting specs."),
]
y = Inches(3.25)
for bold_lead, rest in ai_rows:
    add_dot(s, MARGIN + Inches(0.03), y + Inches(0.11), Inches(0.035), BLUE_VIVID)
    add_runs(s, MARGIN + Inches(0.2), y, Inches(9.5), Inches(0.5),
              [{"text": bold_lead, "font": FONT_BODY, "size": 12.5, "bold": True, "color": WHITE},
               {"text": rest, "font": FONT_BODY, "size": 12.5, "color": RGBColor.from_string("D8DED2")}],
              line_spacing=1.25)
    y = y + Inches(0.5)

note = add_dashed_card(s, MARGIN, y + Inches(0.2), Inches(9.3), Inches(0.95), RGBColor.from_string("6A8F5A"))
add_text(s, MARGIN + Inches(0.25), y + Inches(0.35), Inches(8.8), Inches(0.7),
          "This deck: built with a skill that samples Taxfix’s own rebrand page directly "
          "for its palette, type, and layout rules — the same discipline applied everywhere else.",
          font=FONT_BODY, size=11, color=EG_CALM)
footer(s, 8, TOTAL, RGBColor.from_string("6E7E70"))

# =================================================================
# SLIDE 9 — VALIDATION
# =================================================================
s = add_slide(OW_LIGHT)
eyebrow(s, "08 · VALIDATION", EG_DARK)
add_text(s, MARGIN, Inches(1.15), Inches(8), Inches(1.4), "One metric. One cheap test.",
          font=FONT_HEAD, size=30, color=EG_VERYDARK, line_spacing=1.05)

metrics = [
    ("NORTH STAR", "≥40%", " margin/case", "Once real advisor hours replace the 5-hour placeholder."),
    ("GUARDRAIL", "~2.7x", " LTV:CAC", "Against a 3x bar, at a placeholder €90 CAC. An open risk, stated, not hidden."),
]
mw = Inches(4.6)
x = MARGIN
for label, big, small, desc in metrics:
    add_card(s, x, Inches(2.9), mw, Inches(1.75), WHITE, line_color=OW_CALM)
    add_text(s, x + Inches(0.25), Inches(3.08), mw - Inches(0.5), Inches(0.3), label,
              font=FONT_HEAD, size=10.5, color=RGBColor.from_string("6B6459"))
    add_runs(s, x + Inches(0.25), Inches(3.42), mw - Inches(0.5), Inches(0.55),
              [{"text": big, "font": FONT_HEAD, "size": 26, "color": EG_VERYDARK},
               {"text": small, "font": FONT_BODY, "size": 13, "color": RGBColor.from_string("6B6459")}])
    add_text(s, x + Inches(0.25), Inches(4.02), mw - Inches(0.5), Inches(0.55), desc,
              font=FONT_BODY, size=10.5, color=RGBColor.from_string("6B6459"))
    x = x + mw + Inches(0.3)

bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN, Inches(5.05), Inches(0.04), Inches(0.55))
bar2.fill.solid(); bar2.fill.fore_color.rgb = GOLD_VIVID; bar2.line.fill.background(); bar2.shadow.inherit = False
add_text(s, MARGIN + Inches(0.2), Inches(5.05), Inches(8.5), Inches(0.6),
          "Next: ship the landing page pre-launch and read real conversion and lead "
          "cost — before spending a euro on paid acquisition.",
          font=FONT_BODY, size=13, color=EG_VERYDARK)
footer(s, 9, TOTAL, EG_VERYDARK)

# =================================================================
# SLIDE 10 — CLOSING
# =================================================================
s = add_slide(BLACK)
mark = add_card(s, MARGIN, Inches(0.8), Inches(0.6), Inches(0.6), EG_VIVID, corner=0.28)
mark.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = mark.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "%"; r.font.name = FONT_HEAD; r.font.size = Pt(20); r.font.color.rgb = EG_VERYDARK

add_text(
    s, MARGIN, Inches(1.7), Inches(8.5), Inches(2.8),
    "Small enough that expert-matching doesn’t need millions of cases. "
    "Real enough that the exclusion is a quote off Taxfix’s own product, "
    "not an assumption.",
    font=FONT_HEAD, size=27, color=WHITE, line_spacing=1.15,
)

add_text(
    s, MARGIN, Inches(6.1), Inches(6), Inches(0.4),
    "github.com/hfritz/taxfix-case-study-prototype",
    font=FONT_BODY, size=13, color=EG_VIVID,
)
add_text(
    s, Inches(8.3), Inches(6.0), Inches(4.3), Inches(0.75),
    "Built by Helmut Fritz using AI tools · 2026. Unofficial case-study "
    "prototype — not affiliated with or endorsed by Taxfix.",
    font=FONT_BODY, size=10, color=RGBColor.from_string("A8A8A8"),
)
footer(s, 10, TOTAL, OW_LIGHT)

# ---------------------------------------------------------------
out_path = "/Users/hfritz/code/taxfix/decks/process-case-study-deck.pptx"
prs.save(out_path)
print("Saved:", out_path)
print("Slides:", len(prs.slides._sldIdLst))
